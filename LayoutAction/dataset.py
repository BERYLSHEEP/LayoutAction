
import numpy as np
import torch
from torch.utils.data.dataset import Dataset
from PIL import Image, ImageDraw, ImageOps
import json
import os
from pycocotools.coco import COCO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE

import seaborn as sns
from tqdm import tqdm

def convert_xywh_to_ltrb(bbox):
    xc, yc, w, h = bbox
    x1 = xc - w / 2
    y1 = yc - h / 2
    x2 = xc + w / 2
    y2 = yc + h / 2
    return [x1, y1, x2, y2]

class LexicographicSort():
    def __call__(self, data):
        assert not data.attr['has_canvas_element']
        l, t, _, _ = convert_xywh_to_ltrb(data.x.t())
        _zip = zip(*sorted(enumerate(zip(t, l)), key=lambda c: c[1:]))
        idx = list(list(_zip)[0])
        data.x_orig, data.y_orig = data.x, data.y
        data.x, data.y = data.x[idx], data.y[idx]
        return data

class Padding(object):
    def __init__(self, max_length, vocab_size):
        self.max_length = max_length
        self.bos_token = vocab_size - 3
        self.eos_token = vocab_size - 2
        self.pad_token = vocab_size - 1

    def __call__(self, layout):
        # grab a chunk of (max_length + 1) from the layout

        chunk = torch.zeros(self.max_length+1, dtype=torch.long) + self.pad_token
        # Assume len(item) will always be <= self.max_length:
        chunk[0] = self.bos_token
        chunk[1:len(layout)+1] = layout
        chunk[len(layout)+1] = self.eos_token

        x = chunk[:-1]
        y = chunk[1:]
        return {'x': x, 'y': y}

        # fully reconstruction
        # return {'x': chunk[:-1], 'y': chunk[:-1]}

def get_dataset(name, split, max_length=None, precision=8, is_rela=False):
    if name == "rico":
        if is_rela:
            return MedianExpRelaRicoLayout(split, max_length, precision)
        else:
            return MedianExpRicoLayout(split, max_length, precision)
    elif name == "publaynet":
        return MedianExpPubLayout(split, max_length, precision)
    elif name == "infoppt":
        return MedianExpPPTLayout(split, max_length, precision)
    
    raise NotImplementedError(name)


class BaseDataset(Dataset):
    component_class = []
    _category_id_to_category_name = None
    _json_category_id_to_contiguous_id = None
    _contiguous_id_to_json_id = None
    _colors = None

    def __init__(self, name, split, is_rela=False):
        super().__init__()
        assert split in ['train', 'val', 'test']
        self.is_rela = is_rela
        dir_path = os.path.dirname(os.path.realpath(__file__))
        if is_rela:
            idx = self.processed_file_names.index('{}_rela.pt'.format(split))
        else:
            idx = self.processed_file_names.index('{}.pt'.format(split))
        os.makedirs(os.path.join(dir_path, "preprocess_data", name), exist_ok=True)
        self.data_path = os.path.join(dir_path, "preprocess_data", name, self.processed_file_names[idx])

    @property
    def category_id_to_category_name(self):
        if self._category_id_to_category_name is None:
            self._category_id_to_category_name = {
                self.component_class[key]: key for key in self.component_class.keys()
            }
        return self._category_id_to_category_name

    @property
    def json_category_id_to_contiguous_id(self):
        if self._json_category_id_to_contiguous_id is None:
            self._json_category_id_to_contiguous_id = {
            i: i + self.size + 1 for i in range(self.categories_num)
        }
        return self._json_category_id_to_contiguous_id

    @property
    def contiguous_category_id_to_json_id(self):
        if self._contiguous_id_to_json_id is None:
            self._contiguous_category_id_to_json_id = {
            v: k for k, v in self.json_category_id_to_contiguous_id.items()
        }
        return self._contiguous_id_to_json_id

    @property
    def colors(self):
        if self._colors is None:
            num_colors = self.categories_num
            palette = sns.color_palette(None, num_colors)
            if num_colors > 10:
                palette[10:] = sns.color_palette("husl", num_colors-10)
            self._colors = [[int(x[0]*255), int(x[1]*255), int(x[2]*255)] for x in palette]
        return self._colors

    @property
    def processed_file_names(self):
        if self.is_rela:
            return['train_rela.pt', 'val_rela.pt', 'test.pt']
        else:
            return ['train.pt', 'val.pt', 'test.pt']

    def trim_tokens(self, tokens):
        bos_idx = np.where(tokens == self.bos_token)[0]
        tokens = tokens[bos_idx[0]+1:] if len(bos_idx) > 0 else tokens
        eos_idx = np.where(tokens == self.eos_token)[0]
        tokens = tokens[:eos_idx[0]] if len(eos_idx) > 0 else tokens
        # tokens = tokens[tokens != bos]
        # tokens = tokens[tokens != eos]
        if self.pad_token is not None:
            tokens = tokens[tokens != self.pad_token]
        return tokens

    def get_choice_gt(self, bbox_gt):
        # generate, copy, margin
        copy_label = self.get_copy_gt(bbox_gt)
        # [pre_bbox_num, 4]
        copy_choice = np.sum(copy_label, axis=0, dtype=np.bool_)
        copy_is_zero = copy_choice == np.zeros(copy_choice.shape)
        margin_label, margin_value = self.get_margin_gt(bbox_gt)
        # [pre_bbox_num, 2]
        margin_labelx4 = np.concatenate((margin_label, np.zeros(margin_label.shape)), axis=1)
        margin_choice = np.sum(margin_labelx4, axis=0, dtype=np.bool_) * copy_is_zero
        margin_is_zero = margin_choice == np.zeros(margin_choice.shape)

        generate_choice = copy_is_zero*margin_is_zero

        choice_gt = np.concatenate((copy_choice[:, np.newaxis], margin_choice[:, np.newaxis], generate_choice[:, np.newaxis]), axis=1)
        
        # normalize the label
        # copy_label = torch.nn.functional.normalize(copy_label.to(torch.float32), dim=1, p=1)
        # margin_label = torch.nn.functional.normalize(margin_label.to(torch.float32), dim=1, p=1)
        return choice_gt, copy_label, margin_label, margin_value
        # return : [4, 3]

    def get_copy_gt(self, bbox_gt):
        # bbox_gt: [cur_bbox_max_num, 4]
        bbox_cur_num, _ = bbox_gt.shape
        bbox_g = bbox_gt[-1, :][np.newaxis, :]
        bbox_g = np.repeat(bbox_g, bbox_cur_num-1, 0)
        copy_label = bbox_g == bbox_gt[:-1, :]
        return copy_label

    def get_margin_gt(self, bbox_gt):
        # bbox_gt : [cur_bbox_num, 4]
        bbox_cur_num = bbox_gt.shape[0]
        bbox_g = bbox_gt[-1, :][np.newaxis, :]
        bbox_g = np.repeat(bbox_g, bbox_cur_num-1, 0)
        margin_l = bbox_g == bbox_gt[:-1, :]
        margin_label = np.concatenate((margin_l[:, 1][:, np.newaxis], margin_l[:, 0][:, np.newaxis]), axis=1)
        # margin_value label
        # [cur_bbox_num-1, 2]
        margin_value = bbox_g[:, :2]  - bbox_gt[:-1, :2] - 0.5*bbox_gt[:-1, 2:] - 0.5*bbox_g[:, 2:]
        # if the margin_value < 0: ---> then the corresponding margin label is False.
        is_above_zeros = np.copy(margin_value)
        is_above_zeros[is_above_zeros>=0] = True
        is_above_zeros[is_above_zeros<0] = False

        margin_label = margin_label * is_above_zeros
        
        return margin_label, margin_value

    def save_pt(self, save_path):
        '''
        results = {}
        results["data"] = self.data
        results["iou_data"] = self.iou_data
        '''
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["data"] = self.data
        results["iou_data"] = self.iou_data
        torch.save(results, save_path)

    def load_pt(self, load_path):
        results = torch.load(load_path)
        self.categories_num = results["categories_num"]
        self.max_elements_num = results["max_elements_num"]
        self.data = results["data"]

        if "iou_data" in results:
            print("load iou data")
            self.iou_data = results["iou_data"]

        self.copy_mode = self.size + 1 + self.categories_num
        self.margin_mode = self.copy_mode + 1
        self.generate_mode = self.margin_mode + 1
        self.no_obj_token = self.generate_mode + 1

    def render(self, layout):
        img = Image.new('RGB', (self.W, self.H), color=(255, 255, 255))
        draw = ImageDraw.Draw(img, 'RGBA')
        layout = layout.reshape(-1)
        layout = self.trim_tokens(layout)
        layout = layout[: len(layout) // 13 * 13].reshape(-1, 13)
        box = layout[:, 1:].astype(np.float32)

        option_box = box.reshape(-1, 4, 3) # [box_num, xywh, [option, obj, value]]
        box_num = option_box.shape[0]
        obj_idx = option_box[:, :, 1] - self.no_obj_token
        is_generate = option_box[:, :, 0] == np.ones((box_num, 4))*self.generate_mode
        generate_value = is_generate * option_box[:, :, 2]
        is_copy = option_box[:, :, 0] == np.ones((box_num, 4))*self.copy_mode
        is_margin = option_box[:, :, 0] == np.ones((box_num, 4))*self.margin_mode

        result_box = option_box[0, :, 2][np.newaxis, :]
        for idx in range(1, len(layout)):
            cur_obj_idx = obj_idx[idx].astype(np.int32)
            cur_obj_idx[cur_obj_idx>len(result_box)] = 0
            cur_obj_idx[cur_obj_idx<0] = 0
            cur_copy_value = is_copy[idx] * result_box[-cur_obj_idx, range(4)]
            cur_bbox = generate_value[idx] + cur_copy_value
            margin_xy = result_box[-cur_obj_idx[:2], range(2)] + 0.5*result_box[-cur_obj_idx[:2], [2, 3]] + 0.5*cur_bbox[2:4] + option_box[idx, [0,1], 2]
            cur_margin_value = is_margin[idx] * np.concatenate((margin_xy, np.zeros(margin_xy.shape)), axis=-1)
            cur_bbox += cur_margin_value
            result_box = np.concatenate((result_box, cur_bbox[np.newaxis, :]), axis=0).astype(np.int32)

        # box[:, [0, 1]] = box[:, [0, 1]] / (self.size - 1) * 255
        # box[:, [2, 3]] = box[:, [2, 3]] / self.size * 256
        # box[:, [2, 3]] = box[:, [0, 1]] + box[:, [2, 3]]
        result_box = result_box / (self.size - 1)

        result_box[:, [0, 2]] = result_box[:, [0, 2]] * self.W
        result_box[:, [1, 3]] = result_box[:, [1, 3]] * self.H
        # xywh to ltrb
        x1s = result_box[:, 0] - result_box[:, 2] / 2
        y1s = result_box[:, 1] - result_box[:, 3] / 2
        x2s = result_box[:, 0] + result_box[:, 2] / 2
        y2s = result_box[:, 1] + result_box[:, 3] / 2

        for i in range(len(layout)):
            # x1, y1, x2, y2 = box
            x1, y1, x2, y2 = x1s[i], y1s[i], x2s[i], y2s[i]
            cat = layout[i][0]
            # col = self.colors[cat-self.size] if 0 <= cat-self.size < len(self.colors) else [0, 0, 0]
            col = self.colors[int(cat-self.no_value_token-1)] if 0 <= cat-self.no_value_token-1 < len(self.colors) else [0, 0, 0]
            draw.rectangle([x1, y1, x2, y2],
                           outline=tuple(col) + (200,),
                           fill=tuple(col) + (64,),
                           )

        # Add border around image
        img = ImageOps.expand(img, border=2)
        return img

    def render_normalized_layout(self, layout):
        layout = layout.reshape(-1)
        layout = self.trim_tokens(layout)
        layout = layout[: len(layout) // 13 * 13].reshape(-1, 13)
        if layout.shape[0] == 0:
            return None
        
        box = layout[:, 1:].astype(np.float32)
        option_box = box.reshape(-1, 4, 3) # [box_num, xywh, [option, obj, value]]
        box_num = option_box.shape[0]
        obj_idx = option_box[:, :, 1] - self.no_obj_token
        is_generate = option_box[:, :, 0] == np.ones((box_num, 4))*self.generate_mode
        generate_value = is_generate * option_box[:, :, 2]
        is_copy = option_box[:, :, 0] == np.ones((box_num, 4))*self.copy_mode
        is_margin = option_box[:, :, 0] == np.ones((box_num, 4))*self.margin_mode

        result_box = option_box[0, :, 2][np.newaxis, :]
        for idx in range(1, len(layout)):
            cur_obj_idx = obj_idx[idx].astype(np.int32)
            cur_obj_idx[cur_obj_idx>len(result_box)] = 0
            cur_obj_idx[cur_obj_idx<0] = 0
            cur_copy_value = is_copy[idx] * result_box[-cur_obj_idx, range(4)]
            cur_bbox = generate_value[idx] + cur_copy_value
            margin_xy = result_box[-cur_obj_idx[:2], range(2)] + 0.5*result_box[-cur_obj_idx[:2], [2, 3]] + 0.5*cur_bbox[2:4] + option_box[idx, [0,1], 2]
            cur_margin_value = is_margin[idx] * np.concatenate((margin_xy, np.zeros(margin_xy.shape)), axis=-1)
            cur_bbox += cur_margin_value
            result_box = np.concatenate((result_box, cur_bbox[np.newaxis, :]), axis=0).astype(np.int32)

        # box[:, [0, 1]] = box[:, [0, 1]] / (self.size - 1) * 255
        # box[:, [2, 3]] = box[:, [2, 3]] / self.size * 256
        # box[:, [2, 3]] = box[:, [0, 1]] + box[:, [2, 3]]
        result_box = result_box / (self.size - 1)
        result_box = np.clip(result_box, 0, 1)
        label = layout[:, 0].astype(np.int32) - (self.size+1)
        # label[label>self.categories_num] = 0
        # label[label<0] = 0
        return (result_box, label)

    def calculate_prob(self, layout):
        layout = layout.reshape(-1)
        layout = self.trim_tokens(layout)
        layout = layout[: len(layout) // 13 * 13].reshape(-1, 13)
        
        box = layout[:, 1:].astype(np.float32)
        option_box = box.reshape(-1, 4, 3) # [box_num, xywh, [option, obj, value]]
        box_num = option_box.shape[0]
        is_generate = option_box[:, :, 0] == np.ones((box_num, 4))*self.generate_mode
        is_copy = option_box[:, :, 0] == np.ones((box_num, 4))*self.copy_mode
        is_margin = option_box[:, :2, 0] == np.ones((box_num, 2))*self.margin_mode

        generate_chosen = np.sum(is_generate)
        copy_chosen = np.sum(is_copy)
        margin_chosen = np.sum(is_margin)

        return box_num, generate_chosen, copy_chosen, margin_chosen

    def layout2token(self, layout):
        layout_list = list(map(int, layout))
        
        for idx, token in enumerate(layout_list):
            if isinstance(token, int):
                if token == self.no_value_token:
                    layout_list[idx] = "No_value"
                elif token == self.pad_token:
                    layout_list[idx] = "PAD"
                elif token == self.eos_token:
                    layout_list[idx] = "EOS"
                elif token == self.bos_token:
                    layout_list[idx] = "BOS"
                elif token == self.copy_mode:
                    layout_list[idx] = "COPY"
                elif token == self.margin_mode:
                    layout_list[idx] = "SPACE"
                elif token == self.generate_mode:
                    layout_list[idx] = "GEN"
                elif self.no_value_token < token < self.copy_mode:
                    layout_list[idx] = self.category_id_to_category_name[token - self.no_value_token - 1]
                elif token == self.no_obj_token:
                    if layout_list[idx-1] == "GEN":
                        layout_list[idx] = "NoObj"
                    else:
                        layout_list[idx] = 1
                elif self.no_obj_token < token < self.pad_token:
                    layout_list[idx] = token - self.no_obj_token
            
        return layout_list

    def quantize_box(self, boxes, width, height):
        # range of xy is [0, large_side-1]
        # range of wh is [1, large_side]
        # bring xywh to [0, 1]
        
        boxes[:, [2, 3]] = boxes[:, [2, 3]] - 1
        boxes[:, [0, 2]] = boxes[:, [0, 2]] / (width - 1)
        boxes[:, [1, 3]] = boxes[:, [1, 3]] / (height - 1)
        boxes = np.clip(boxes, 0, 1)

        # next take xywh to [0, size-1]
        boxes = (boxes * (self.size - 1)).round()
        # round: 四舍五入

        return boxes.astype(np.int32)

    def __len__(self):
        return len(self.data)

    def __getitem__(self, idx):
        # grab a chunk of (block_size + 1) tokens from the data
        layout = torch.tensor(self.data[idx], dtype=torch.long)
        layout = self.transform(layout)
        return layout['x'], layout['y'] 

class MedianExpRelaRicoLayout(BaseDataset):
    component_class = {'Toolbar':0, 'Image':1, 'Text':2, 'Icon':3, 'Text Button':4, 'Input':5,
        'List Item': 6, 'Advertisement': 7, 'Pager Indicator':8, 'Web View':9, 'Background Image':10,
        'Drawer':11, 'Modal':12}

    def __init__(self, split, max_length=None, precision=8):    
        super().__init__('rico', split, is_rela=True)

        self.W = 256
        self.H = 256
        self.size = pow(2, precision)
        self.no_value_token = self.size

        if os.path.exists(self.data_path):
            print("load dataset.")
            self.load_pt(self.data_path)
        else:
            data_dir = f"./datasets/rico/semantic_annotations"
            dirs = os.listdir(data_dir)
            self.max_elements_num = 9
            self.categories_num = len(self.component_class.keys())

            self.copy_mode = self.size + 1 + self.categories_num
            self.margin_mode = self.copy_mode + 1
            self.generate_mode = self.margin_mode + 1
            self.option_id = np.array([self.copy_mode, self.margin_mode, self.generate_mode])

            self.no_obj_token = self.generate_mode + 1
            self.obj_id_to_contiguous_id = {
                i: i + self.generate_mode + 1 for i in range(self.max_elements_num)
            }
            
            self.data = []
            self.rela_dict = {}
            self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
            bbox_idx = 0

            for file in tqdm(dirs, total=len(dirs)):
                if file.split(".")[-1] == "json":
                    file_path = os.path.join(data_dir, file)
                    with open(file_path, encoding='utf-8') as f:
                        json_file = json.load(f)

                    canvas = json_file["bounds"]
                    W, H = float(canvas[2]-canvas[0]), float(canvas[3]-canvas[1])
                    if canvas[0]!= 0 or canvas[1]!= 0 or W <= 1000:
                        continue
                    elements = self.get_all_element(json_file, [])
                    elements = list(filter(lambda e: e["componentLabel"] in self.component_class, elements))
                    
                    if len(elements) == 0 or len(elements)> self.max_elements_num:
                        continue
                    
                    ann_box = []
                    ann_cat = []

                    for ele in elements:
                        [x_l, y_t, x_r, y_b] = ele["bounds"]
                        xc = (x_l + x_r) / 2.
                        yc = (y_t + y_b) / 2.
                        w = x_r - x_l
                        h = y_b - y_t

                        if w<0 or h<0:
                            continue
                        ann_box.append([xc, yc, w, h])
                        ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[ele["componentLabel"]]])

                    # Sort boxes

                    ann_box = np.array(ann_box)
                    # Discretize boxes
                    ann_box = self.quantize_box(ann_box, W, H)

                    ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                    # Sort by ann_box[:, 1], then by ann_box[:, 0]
                    ann_box = ann_box[ind]
                    
                    ann_cat = np.array(ann_cat)
                    ann_cat = ann_cat[ind]

                    self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                    self.iou_data["file_idx"].append(int(file.split(".")[0]))
                    self.iou_data["file2bboxidx"][int(file.split(".")[0])] = bbox_idx

                    # xywh to Option Obj Value
                    # the first element must from generate option
                    ann_option = [[self.generate_mode]*4]
                    ann_obj = [[self.no_obj_token]*4]
                    ann_value = [ann_box[0]]
                    
                    for ele_idx in range(1, len(ann_box)):
                        choice_gt, copy_label, margin_label, margin_value = self.get_choice_gt(ann_box[:ele_idx+1])
                        # option :  choice_gt [bs, 4, 3]  -->  ann_option [bs, 4]
                        option_idx = np.argmax(choice_gt, axis=-1)
                        ann_option.append(self.option_id[option_idx])

                        # obj: copy_label [pre_num, 4]  margin_label [pre_num, 4]   ---> ann_obj [4]
                        copy_obj_flip_idx = np.argmax(np.flip(copy_label, 0), axis=0)
                        # if the idx is 1, indicates the last (1 + idx) = 2 obj is selected.
                        margin_label_4dim = np.concatenate((margin_label, np.zeros(margin_label.shape)), axis=1)
                        margin_obj_flip_idx = np.argmax(np.flip(margin_label_4dim, 0), axis=0) 

                        copy_obj_hit_idx = (copy_obj_flip_idx + 1) * choice_gt[:, 0]
                        margin_obj_hit_idx = (margin_obj_flip_idx + 1) * choice_gt[:, 1]
                        obj_idx = copy_obj_hit_idx + margin_obj_hit_idx + self.no_obj_token
                        ann_obj.append(obj_idx)

                        copy_hit_value = choice_gt[:, 0] * self.no_value_token
                        margin_value_x4 = np.concatenate((margin_value[-margin_obj_hit_idx[:2], range(2)][np.newaxis, :], np.zeros((1, 2))), axis=1) 
                        margin_hit_value = choice_gt[:, 1] * margin_value_x4
                        generate_hit_value = choice_gt[:, 2] * ann_box[ele_idx]
                        value_idx = copy_hit_value + margin_hit_value + generate_hit_value
                        ann_value.append(value_idx[0].round().astype(np.int32))

                        for geo in range(4):
                            hit_global = False
                            if choice_gt[geo, 0] == 1:
                                oidx = ele_idx - copy_obj_flip_idx[geo] -1
                                key = list(ann_box[oidx]) + [ann_cat[oidx], geo, 0, 0] # 0: copy mode
                                hit_global = True
                            elif choice_gt[geo, 1] == 1:
                                oidx = ele_idx - margin_obj_flip_idx[geo] - 1
                                mvalue = margin_value_x4[0, geo]
                                key = list(ann_box[oidx]) + [ann_cat[oidx], geo, 1, mvalue] # 1: margin mode
                                hit_global = True
                            if hit_global:
                                value = [bbox_idx, ele_idx, geo]
                                if str(key) in self.rela_dict:
                                    self.rela_dict[str(key)].append(value)
                                else:
                                    self.rela_dict[str(key)] = [value]

                    ann_option = np.array(ann_option)
                    ann_obj = np.array(ann_obj)
                    ann_value = np.array(ann_value)

                    OOV_idx = np.concatenate([ann_option[:, :, np.newaxis], ann_obj[:, :, np.newaxis], ann_value[:, :, np.newaxis]], axis=-1)
                    layout = np.concatenate([ann_cat[:, np.newaxis], OOV_idx.reshape(-1, 12)], axis=-1)
                    
                    # Flatten and add to the dataset
                    bbox_idx += 1
                    self.data.append(layout.reshape(-1))

            inner_batch_size = 32
            self.inbatch_data = []
            for key in self.rela_dict.keys():
                rela_all = self.rela_dict[key]
                if len(rela_all) < 2:
                    continue
                rela_layers = []
                rela_labels_idx = []
                for i, (data_idx, ele_idx, geo_idx) in enumerate(rela_all):
                    rela_layer = self.data[data_idx][np.newaxis, :]
                    rela_label_idx = ele_idx*13 + 1 + geo_idx*3 + 1
                    rela_layers.append(rela_layer)
                    rela_labels_idx.append(rela_label_idx)
                    if (i+1) % inner_batch_size == 0 or (i+1) == len(rela_all):
                        self.inbatch_data.append([rela_layers, rela_labels_idx])
                        rela_layers = []
                        rela_labels_idx = []

            self.save_pt(self.data_path)  

        self.max_length = max_length
        if self.max_length is None:
            max_length_cxywh = max([len(x) for x in self.data])
            self.max_elements_num =  int(max_length_cxywh/ 13)
            # category + (option + obj + value) * (xywh) = 13
            self.max_length = max_length_cxywh + 2  # bos, eos tokens
        
        self.vocab_size = self.size + 1 + self.categories_num + 3 + 1 + self.max_elements_num + 3  
        # size, no value, category num, copy, margin, generate, no obj, obj idx, bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1
        self.transform = Padding(self.max_length, self.vocab_size)

    def __len__(self):
        return len(self.inbatch_data)

    def __getitem__(self, idx):
        # grab a chunk of (block_size + 1) tokens from the data
        in_data = self.inbatch_data[idx][0]
        batch_x = []
        batch_y = []
        for i in range(len(in_data)):
            layout = torch.tensor(in_data[i].squeeze(0), dtype=torch.long)
            layout = self.transform(layout)
            batch_x.append(layout['x'].unsqueeze(0))
            batch_y.append(layout['y'].unsqueeze(0))
        batch_x = torch.cat(batch_x, dim=0)
        batch_y = torch.cat(batch_y, dim =0)
        batch_rela_y = torch.tensor(self.inbatch_data[idx][1], dtype=torch.long)
        return batch_x, batch_y, batch_rela_y

    def load_pt(self, load_path):
        results = torch.load(load_path)
        self.categories_num = results["categories_num"]
        self.max_elements_num = results["max_elements_num"]
        self.inbatch_data = results["inbatch_data"]
        self.data = results["data"]

        if "iou_data" in results:
            print("load iou data")
            self.iou_data = results["iou_data"]

        self.copy_mode = self.size + 1 + self.categories_num
        self.margin_mode = self.copy_mode + 1
        self.generate_mode = self.margin_mode + 1
        self.no_obj_token = self.generate_mode + 1

    def save_pt(self, save_path):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data))
        s = [int(N * .85), int(N * .90)]
        results["inbatch_data"] = self.inbatch_data[:s[0]]
        results["data"] = self.data[:s[0]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[0]))
        results["inbatch_data"] = self.inbatch_data[s[0]:s[1]]
        results["data"] = self.data[s[0]:s[1]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[1]))
        results["inbatch_data"] = self.inbatch_data[s[1]:]
        results["data"] = self.data[s[1]:]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[2]))

    def get_all_element(self, p_dic, elements):
        if "children" in p_dic:
            for i in range(len(p_dic["children"])):
                cur_child = p_dic["children"][i]
                elements.append(cur_child)
                elements = self.get_all_element(cur_child, elements)
        return elements       

class MedianExpRicoLayout(BaseDataset):
    # component_class = {'Text':0, 'Icon':1, 'Image':2, 'Text Button':3, 'Toolbar':4, 'List Item':5, 'Web View':6, 
    # 'Advertisement':7, 'Input':8, 'Drawer':9, 'Background Image':10, 'Card':11, 'Multi-Tab':12, 'Modal':13, 
    # 'Pager Indicator':14, 'Radio Button':15, 'On/Off Switch':16, 'Slider':17, 'Checkbox':18, 'Map View':19,
    # 'Button Bar':20, 'Video':21, 'Bottom Navigation':22, 'Date Picker':23, 'Number Stepper':24}
    component_class = {'Toolbar':0, 'Image':1, 'Text':2, 'Icon':3, 'Text Button':4, 'Input':5,
        'List Item': 6, 'Advertisement': 7, 'Pager Indicator':8, 'Web View':9, 'Background Image':10,
        'Drawer':11, 'Modal':12}

    def __init__(self, split, max_length=None, precision=8):    
        super().__init__('rico', split)

        self.W = 256
        self.H = 256
        self.size = pow(2, precision)
        self.no_value_token = self.size

        if os.path.exists(self.data_path):
            print("load dataset.")
            self.load_pt(self.data_path)
        else:
            data_dir = f"./datasets/rico/semantic_annotations"
            dirs = os.listdir(data_dir)
            self.max_elements_num = 9
            self.categories_num = len(self.component_class.keys())

            self.copy_mode = self.size + 1 + self.categories_num
            self.margin_mode = self.copy_mode + 1
            self.generate_mode = self.margin_mode + 1
            self.option_id = np.array([self.copy_mode, self.margin_mode, self.generate_mode])

            self.no_obj_token = self.generate_mode + 1
            self.obj_id_to_contiguous_id = {
                i: i + self.generate_mode + 1 for i in range(self.max_elements_num)
            }
            
            self.data = []

            self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
            bbox_idx = 0

            for file in tqdm(dirs, total=len(dirs)):
                if file.split(".")[-1] == "json":
                    file_path = os.path.join(data_dir, file)
                    with open(file_path, encoding='utf-8') as f:
                        json_file = json.load(f)

                    canvas = json_file["bounds"]
                    W, H = float(canvas[2]-canvas[0]), float(canvas[3]-canvas[1])
                    if canvas[0]!= 0 or canvas[1]!= 0 or W <= 1000:
                        continue
                    elements = self.get_all_element(json_file, [])
                    elements = list(filter(lambda e: e["componentLabel"] in self.component_class, elements))
                    
                    if len(elements) == 0 or len(elements)> self.max_elements_num:
                        continue
                    
                    ann_box = []
                    ann_cat = []

                    for ele in elements:
                        [x_l, y_t, x_r, y_b] = ele["bounds"]
                        xc = (x_l + x_r) / 2.
                        yc = (y_t + y_b) / 2.
                        w = x_r - x_l
                        h = y_b - y_t

                        if w<0 or h<0:
                            continue
                        ann_box.append([xc, yc, w, h])
                        ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[ele["componentLabel"]]])

                    # Sort boxes

                    ann_box = np.array(ann_box)
                    # Discretize boxes
                    ann_box = self.quantize_box(ann_box, W, H)

                    ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                    # Sort by ann_box[:, 1], then by ann_box[:, 0]
                    ann_box = ann_box[ind]
                    
                    ann_cat = np.array(ann_cat)
                    ann_cat = ann_cat[ind]

                    self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                    self.iou_data["file_idx"].append(int(file.split(".")[0]))
                    self.iou_data["file2bboxidx"][int(file.split(".")[0])] = bbox_idx

                    bbox_idx += 1

                    # xywh to Option Obj Value
                    # the first element must from generate option
                    ann_option = [[self.generate_mode]*4]
                    ann_obj = [[self.no_obj_token]*4]
                    ann_value = [ann_box[0]]
                    
                    for ele_idx in range(1, len(ann_box)):
                        choice_gt, copy_label, margin_label, margin_value = self.get_choice_gt(ann_box[:ele_idx+1])
                        # option :  choice_gt [bs, 4, 3]  -->  ann_option [bs, 4]
                        option_idx = np.argmax(choice_gt, axis=-1)
                        ann_option.append(self.option_id[option_idx])

                        # obj: copy_label [pre_num, 4]  margin_label [pre_num, 4]   ---> ann_obj [4]
                        copy_obj_flip_idx = np.argmax(np.flip(copy_label, 0), axis=0)
                        # if the idx is 1, indicates the last (1 + idx) = 2 obj is selected.
                        margin_label_4dim = np.concatenate((margin_label, np.zeros(margin_label.shape)), axis=1)
                        margin_obj_flip_idx = np.argmax(np.flip(margin_label_4dim, 0), axis=0) 

                        copy_obj_hit_idx = (copy_obj_flip_idx + 1) * choice_gt[:, 0]
                        margin_obj_hit_idx = (margin_obj_flip_idx + 1) * choice_gt[:, 1]
                        obj_idx = copy_obj_hit_idx + margin_obj_hit_idx + self.no_obj_token
                        ann_obj.append(obj_idx)

                        copy_hit_value = choice_gt[:, 0] * self.no_value_token
                        margin_value_x4 = np.concatenate((margin_value[-margin_obj_hit_idx[:2], range(2)][np.newaxis, :], np.zeros((1, 2))), axis=1) 
                        margin_hit_value = choice_gt[:, 1] * margin_value_x4
                        generate_hit_value = choice_gt[:, 2] * ann_box[ele_idx]
                        value_idx = copy_hit_value + margin_hit_value + generate_hit_value
                        ann_value.append(value_idx[0].round().astype(np.int32))

                    ann_option = np.array(ann_option)
                    ann_obj = np.array(ann_obj)
                    ann_value = np.array(ann_value)

                    OOV_idx = np.concatenate([ann_option[:, :, np.newaxis], ann_obj[:, :, np.newaxis], ann_value[:, :, np.newaxis]], axis=-1)
                    layout = np.concatenate([ann_cat[:, np.newaxis], OOV_idx.reshape(-1, 12)], axis=-1)
                    
                    # Flatten and add to the dataset
                    self.data.append(layout.reshape(-1))

            self.save_pt(self.data_path)  

        self.max_length = max_length
        if self.max_length is None:
            max_length_cxywh = max([len(x) for x in self.data])
            self.max_elements_num =  int(max_length_cxywh/ 13)
            # category + (option + obj + value) * (xywh) = 13
            self.max_length = max_length_cxywh + 2  # bos, eos tokens
        
        self.vocab_size = self.size + 1 + self.categories_num + 3 + 1 + self.max_elements_num + 3  
        # size, no value, category num, copy, margin, generate, no obj, obj idx, bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1
        self.transform = Padding(self.max_length, self.vocab_size)

    def save_pt(self, save_path):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data))
        s = [int(N * .85), int(N * .90)]
        results["data"] = self.data[:s[0]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[0]))
        results["data"] = self.data[s[0]:s[1]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[1]))
        results["data"] = self.data[s[1]:]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[2]))

    def get_all_element(self, p_dic, elements):
        if "children" in p_dic:
            for i in range(len(p_dic["children"])):
                cur_child = p_dic["children"][i]
                elements.append(cur_child)
                elements = self.get_all_element(cur_child, elements)
        return elements
   
class MedianExpPubLayout(BaseDataset):
    component_class = {'text': 0, 'title':1, 'list':2, 'table':3, 'figure':4}
        
    def __init__(self, split, max_length=None, precision=8):
        super().__init__('publaynet', split)

        self.W = 256
        self.H = 256
        self.size = pow(2, precision)
        self.no_value_token = self.size

        self.max_elements_num = 9
        self.categories_num = len(self.component_class.keys())
        
        self.copy_mode = self.size + 1 + self.categories_num
        self.margin_mode = self.copy_mode + 1
        self.generate_mode = self.margin_mode + 1
        self.option_id = np.array([self.copy_mode, self.margin_mode, self.generate_mode])

        self.no_obj_token = self.generate_mode + 1
        self.obj_id_to_contiguous_id = {
            i: i + self.generate_mode + 1 for i in range(self.max_elements_num)
        }
        
        self.data = []
        self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
        bbox_idx = 0
        
        if os.path.exists(self.data_path):
            print("load dataset.")
            self.load_pt(self.data_path)
        else:
            ori_split = "train" if split in ["train", "val"] else "val"
            data_path = f"./datasets/publaynet/{ori_split}.json"
            coco = COCO(data_path)
            for img_id in sorted(coco.getImgIds()):
                ann_img = coco.loadImgs(img_id)
                W = float(ann_img[0]['width'])
                H = float(ann_img[0]['height'])
                name = ann_img[0]['file_name']
                if H < W:
                    continue

                def is_valid(element):
                    x1, y1, width, height = element['bbox']
                    x2, y2 = x1 + width, y1 + height
                    if x1 < 0 or y1 < 0 or W < x2 or H < y2:
                        return False
                    if x2 <= x1 or y2 <= y1:
                        return False
                    if width <=0 or height <=0:
                        return False
                    return True

                elements = coco.loadAnns(coco.getAnnIds(imgIds=[img_id]))
                elements = list(filter(is_valid, elements))
                
                N = len(elements)
                if N == 0 or 9 < N:
                    continue

                ann_box = []
                ann_cat = []

                for element in elements:
                    # bbox
                    x1, y1, width, height = element['bbox']
                    xc = x1 + width / 2.
                    yc = y1 + height / 2.
                    b = [xc , yc , width, height]
                    ann_box.append(b)

                    # label
                    l = coco.cats[element['category_id']]['name']
                    ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[l]])

                # Sort boxes
                ann_box = np.array(ann_box)
                # Discretize boxes
                ann_box = self.quantize_box(ann_box, W, H)

                ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                # Sort by ann_box[:, 1], then by ann_box[:, 0]
                ann_box = ann_box[ind]
                
                ann_cat = np.array(ann_cat)
                ann_cat = ann_cat[ind]

                self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                self.iou_data["file_idx"].append(img_id)
                self.iou_data["file2bboxidx"][img_id] = bbox_idx

                bbox_idx += 1

                # xywh to Option Obj Value
                # the first element must from generate option
                ann_option = [[self.generate_mode]*4]
                ann_obj = [[self.no_obj_token]*4]
                ann_value = [ann_box[0]]
                
                for ele_idx in range(1, len(ann_box)):
                    choice_gt, copy_label, margin_label, margin_value = self.get_choice_gt(ann_box[:ele_idx+1])
                    # option :  choice_gt [bs, 4, 3]  -->  ann_option [bs, 4]
                    option_idx = np.argmax(choice_gt, axis=-1)
                    ann_option.append(self.option_id[option_idx])

                    # obj: copy_label [pre_num, 4]  margin_label [pre_num, 4]   ---> ann_obj [4]
                    copy_obj_flip_idx = np.argmax(np.flip(copy_label, 0), axis=0)
                    # if the idx is 1, indicates the last (1 + idx) = 2 obj is selected.
                    margin_label_4dim = np.concatenate((margin_label, np.zeros(margin_label.shape)), axis=1)
                    margin_obj_flip_idx = np.argmax(np.flip(margin_label_4dim, 0), axis=0) 

                    copy_obj_hit_idx = (copy_obj_flip_idx + 1) * choice_gt[:, 0]
                    margin_obj_hit_idx = (margin_obj_flip_idx + 1) * choice_gt[:, 1]
                    obj_idx = copy_obj_hit_idx + margin_obj_hit_idx + self.no_obj_token
                    ann_obj.append(obj_idx)

                    copy_hit_value = choice_gt[:, 0] * self.no_value_token
                    margin_value_x4 = np.concatenate((margin_value[-margin_obj_hit_idx[:2], range(2)][np.newaxis, :], np.zeros((1, 2))), axis=1) 
                    margin_hit_value = choice_gt[:, 1] * margin_value_x4
                    generate_hit_value = choice_gt[:, 2] * ann_box[ele_idx]
                    value_idx = copy_hit_value + margin_hit_value + generate_hit_value
                    ann_value.append(value_idx[0].round().astype(np.int32))

                ann_option = np.array(ann_option)
                ann_obj = np.array(ann_obj)
                ann_value = np.array(ann_value)

                OOV_idx = np.concatenate([ann_option[:, :, np.newaxis], ann_obj[:, :, np.newaxis], ann_value[:, :, np.newaxis]], axis=-1)
                layout = np.concatenate([ann_cat[:, np.newaxis], OOV_idx.reshape(-1, 12)], axis=-1)
                
                # Flatten and add to the dataset
                self.data.append(layout.reshape(-1))
       
            self.save_pt(self.data_path, split)

        self.max_length = max_length
        if self.max_length is None:
            max_length_cxywh = max([len(x) for x in self.data])
            self.max_elements_num =  int(max_length_cxywh/ 13)
            # category + (option + obj + value) * (xywh) = 13
            self.max_length = max_length_cxywh + 2  # bos, eos tokens
        
        self.vocab_size = self.size + 1 + self.categories_num + 3 + 1 + self.max_elements_num + 3  
        # little bug here: the vocab_size should be the following definition since only (self.max_elements_num - 1) previous object can be selected.
        # self.vocab_size = self.size + 1 + self.categories_num + 3 + self.max_elements_num + 3  
        # size, no value, category num, copy, margin, generate, no obj, obj idx, bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1
        self.transform = Padding(self.max_length, self.vocab_size)
        
    def save_pt(self, save_path, split):
        if split == "test":
            return super().save_pt(save_path)

        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data)*0.95)
        if split == "train":
            results["data"] = self.data[:N]
        else:
            results["data"] = self.data[N:]
        torch.save(results, save_path)

class MedianExpPPTLayout(BaseDataset):
    component_class = {"TEXT_BOX": 0, "PICTURE": 1, "CHART":2, "TABLE":3, "TITLE":4, "SUBTITLE":5}   
        
    def __init__(self, split, max_length=None, precision=8):
        super().__init__('infoppt', split)

        SHAPE_TYPE_DICT = {}
        PLACEHOLDER_DICT = {}
        AUTO_SHAPE_DICT = {}
        for idx, name in enumerate(dir(MSO_SHAPE_TYPE), start=1):
            if idx > 26:
                break
            tmp_num = getattr(MSO_SHAPE_TYPE, name)
            SHAPE_TYPE_DICT[tmp_num] = name
        for idx, name in enumerate(dir(PP_PLACEHOLDER), start=1):
            if idx > 22:
                break
            tmp_num = getattr(PP_PLACEHOLDER, name)
            PLACEHOLDER_DICT[tmp_num] = name
        for idx, name in enumerate(dir(MSO_SHAPE), start=1):
            if idx > 184:
                break
            tmp_num = getattr(MSO_SHAPE, name)
            AUTO_SHAPE_DICT[tmp_num] = name

        self.categories_num = len(self.component_class.keys())

        # shape.picture + placeholder.picture
        rectangle_class = ['RECTANGLE', 'ROUNDED_RECTANGLE']

        self.W = 256
        self.H = 256
        self.size = pow(2, precision)
        self.no_value_token = self.size

        self.max_elements_num = 20
        self.categories_num = len(self.component_class.keys())

        self.copy_mode = self.size + 1 + self.categories_num
        self.margin_mode = self.copy_mode + 1
        self.generate_mode = self.margin_mode + 1
        self.option_id = np.array([self.copy_mode, self.margin_mode, self.generate_mode])

        # here self.obj_id_to_contiguous_id[0] = self.no_obj_token
        self.obj_id_to_contiguous_id = {
            i: i + self.generate_mode + 1 for i in range(self.max_elements_num)
        }
        self.no_obj_token = self.generate_mode + 1
        
        if os.path.exists(self.data_path):
            print("load InfoPPT dataset.")
            self.load_pt(self.data_path)
        else:
            data_dir = f"./datasets/infoppt"
            dirs = os.listdir(data_dir)
            self.data = []
            self.iou_data = {"bbox":[], "file_idx":[], "file2bboxidx":{}}
            bbox_idx = 0
            for file in tqdm(dirs, total=len(dirs)):        
                if file.split(".")[-1] == "pptx":
                    file_path = os.path.join(data_dir, file)
                    prs = Presentation(file_path)

                    H, W = prs.slide_height, prs.slide_width   
                    for index, slide in enumerate(prs.slides, start=1):
                        ann_box = []
                        ann_cat = []
                        for shape in slide.shapes:
                            try:
                                # shape.left, shape.top, shape.width, shape.height
                                xc = shape.left + shape.width*0.5
                                yc = shape.top + shape.height*0.5

                                if shape.width <= 0 or shape.height <= 0:
                                    continue

                                if shape.is_placeholder:
                                    cur_class = PLACEHOLDER_DICT[shape.placeholder_format.type]
                                    if cur_class not in ["PICTURE", "TITLE", "SUBTITLE"]:
                                        continue
                                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                    cur_class = AUTO_SHAPE_DICT[shape.auto_shape_type]
                                    if cur_class in rectangle_class and shape.has_text_frame:
                                        if len(shape.text) > 0:
                                            cur_class = "TEXT_BOX"
                                        else:
                                            continue
                                    else:
                                        continue
                                elif SHAPE_TYPE_DICT[shape.shape_type] in self.component_class:
                                    cur_class = SHAPE_TYPE_DICT[shape.shape_type]
                                    if cur_class == "TEXT_BOX" and shape.has_text_frame:
                                        if len(shape.text) == 0:
                                            continue
                                else:
                                    # filter
                                    continue                               

                                ann_box.append([xc, yc, shape.width, shape.height])
                                ann_cat.append(self.json_category_id_to_contiguous_id[self.component_class[cur_class]])
                            except:
                                pass
                            
                        if len(ann_cat) > self.max_elements_num or len(ann_cat) <= 3:
                            continue

                        ann_box = np.array(ann_box)

                        # Sort boxes
                        # Discretize boxes
                        ann_box = self.quantize_box(ann_box, W, H)
                        ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                        # Sort by ann_box[:, 1], then by ann_box[:, 0]
                        ann_box = ann_box[ind]
                        
                        ann_cat = np.array(ann_cat)
                        
                        ann_cat = ann_cat[ind]

                        file_name = file.split(".")[0]
                        self.iou_data["bbox"].append(ann_box[np.newaxis, :]/(self.size - 1))
                        self.iou_data["file_idx"].append(f"{file_name}_{index}")
                        self.iou_data["file2bboxidx"][f"{file_name}_{index}"] = bbox_idx

                        bbox_idx += 1
                        

                        # xywh to Option Obj Value
                        # the first element must from generate option
                        ann_option = [[self.generate_mode]*4]
                        ann_obj = [[self.no_obj_token]*4]
                        ann_value = [ann_box[0]]
                        
                        for ele_idx in range(1, len(ann_box)):
                            choice_gt, copy_label, margin_label, margin_value = self.get_choice_gt(ann_box[:ele_idx+1])
                            # option :  choice_gt [bs, 4, 3]  -->  ann_option [bs, 4]
                            option_idx = np.argmax(choice_gt, axis=-1)
                            ann_option.append(self.option_id[option_idx])

                            # obj: copy_label [pre_num, 4]  margin_label [pre_num, 4]   ---> ann_obj [4]
                            copy_obj_flip_idx = np.argmax(np.flip(copy_label, 0), axis=0)
                            # if the idx is 1, indicates the last (1 + idx) = 2 obj is selected.
                            margin_label_4dim = np.concatenate((margin_label, np.zeros(margin_label.shape)), axis=1)
                            margin_obj_flip_idx = np.argmax(np.flip(margin_label_4dim, 0), axis=0) 

                            copy_obj_hit_idx = (copy_obj_flip_idx + 1) * choice_gt[:, 0]
                            margin_obj_hit_idx = (margin_obj_flip_idx + 1) * choice_gt[:, 1]
                            obj_idx = copy_obj_hit_idx + margin_obj_hit_idx + self.no_obj_token
                            ann_obj.append(obj_idx)

                            copy_hit_value = choice_gt[:, 0] * self.no_value_token
                            margin_value_x4 = np.concatenate((margin_value[-margin_obj_hit_idx[:2], range(2)][np.newaxis, :], np.zeros((1, 2))), axis=1) 
                            margin_hit_value = choice_gt[:, 1] * margin_value_x4
                            generate_hit_value = choice_gt[:, 2] * ann_box[ele_idx]
                            value_idx = copy_hit_value + margin_hit_value + generate_hit_value
                            ann_value.append(value_idx[0].round().astype(np.int32))

                        ann_option = np.array(ann_option)
                        ann_obj = np.array(ann_obj)
                        ann_value = np.array(ann_value)

                        OOV_idx = np.concatenate([ann_option[:, :, np.newaxis], ann_obj[:, :, np.newaxis], ann_value[:, :, np.newaxis]], axis=-1)
                        layout = np.concatenate([ann_cat[:, np.newaxis], OOV_idx.reshape(-1, 12)], axis=-1)
                        
                        # Flatten and add to the dataset
                        self.data.append(layout.reshape(-1))

            self.save_pt(self.data_path)

        self.max_length = max_length
        if self.max_length is None:
            max_length_cxywh = max([len(x) for x in self.data])
            self.max_elements_num =  int(max_length_cxywh/ 13)
            # category + (option + obj + value) * (xywh) = 13
            self.max_length = max_length_cxywh + 2  # bos, eos tokens
        
        self.vocab_size = self.size + 1 + self.categories_num + 3 + 1 + self.max_elements_num + 3  
        # size, no value, category num, copy, margin, generate, no obj, obj idx, bos, eos, pad tokens
        self.bos_token = self.vocab_size - 3
        self.eos_token = self.vocab_size - 2
        self.pad_token = self.vocab_size - 1
        self.transform = Padding(self.max_length, self.vocab_size)

    def save_pt(self, save_path):
        results = {}
        results["categories_num"] = self.categories_num
        results["max_elements_num"] = self.max_elements_num
        results["iou_data"] = self.iou_data
        N = int(len(self.data))
        s = [int(N * .85), int(N * .90)]
        results["data"] = self.data[:s[0]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[0]))
        results["data"] = self.data[s[0]:s[1]]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[1]))
        results["data"] = self.data[s[1]:]
        torch.save(results, os.path.join(os.path.dirname(save_path), self.processed_file_names[2]))   