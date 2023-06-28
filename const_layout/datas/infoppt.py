from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE
import torch
from torch_geometric.data import Data
import os
import numpy as np
from datas.base import BaseDataset


SHAPE_TYPE_DICT = {}
PLACEHOLDER_DICT = {}
AUTO_SHAPE_DICT = {}
for idx, name in enumerate(dir(MSO_SHAPE_TYPE), start=1):
    if idx > 26:
        break
    tmp_num = getattr(MSO_SHAPE_TYPE, name)
    SHAPE_TYPE_DICT[tmp_num] = name
for idx, name in enumerate(dir(PP_PLACEHOLDER), start=1):
    if idx > 22:
        break
    tmp_num = getattr(PP_PLACEHOLDER, name)
    PLACEHOLDER_DICT[tmp_num] = name
for idx, name in enumerate(dir(MSO_SHAPE), start=1):
    if idx > 184:
        break
    tmp_num = getattr(MSO_SHAPE, name)
    AUTO_SHAPE_DICT[tmp_num] = name

class PPTNet(BaseDataset):
    labels = [
        "TEXT_BOX",
        "PICTURE",
        "CHART",
        "TABLE",
        "TITLE",
        "SUBTITLE"
    ]
    def __init__(self, split='train', transform=None):
        super().__init__('ppt', split, transform)

    def download(self):
        super().download()

    def process(self):
        ppt_class = {"TEXT_BOX": 0, "PICTURE": 1, "CHART":2, "TABLE":3, "TITLE":4, "SUBTITLE":5}  
        component_dim = len(ppt_class.keys())
        # shape.picture + placeholder.picture
        rectangle_class = ['RECTANGLE', 'ROUNDED_RECTANGLE']

        raw_dir = os.path.join(self.raw_dir, 'infoppt')
        for file in os.listdir(raw_dir):
            data_list = []

            if file.split(".")[-1] == "pptx":
                file_path = os.path.join(raw_dir, file)
                prs = Presentation(file_path)
                
                H, W = prs.slide_height, prs.slide_width   
                max_element_num = 20
                for index, slide in enumerate(prs.slides, start=1):
                    ann_box = []
                    ann_cat = []

                    if int(file.split(".")[0]) == 2 and index == 15:
                        print("2_15")
                    for shape in slide.shapes:
                        try:
                            # shape.left, shape.top, shape.width, shape.height
                            xc = shape.left + shape.width*0.5
                            yc = shape.top + shape.height*0.5

                            if shape.width <= 0 or shape.height <= 0:
                                continue
                            
                            if shape.is_placeholder:
                                cur_class = PLACEHOLDER_DICT[shape.placeholder_format.type]
                                if cur_class not in ["PICTURE", "TITLE", "SUBTITLE"]:
                                    continue
                            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                                cur_class = AUTO_SHAPE_DICT[shape.auto_shape_type]
                                if cur_class in rectangle_class and shape.has_text_frame:
                                    if len(shape.text) > 0:
                                        cur_class = "TEXT_BOX"
                                    else:
                                        continue
                                else:
                                    continue
                            elif SHAPE_TYPE_DICT[shape.shape_type] in ppt_class:
                                cur_class = SHAPE_TYPE_DICT[shape.shape_type]
                                if cur_class == "TEXT_BOX" and shape.has_text_frame:
                                    if len(shape.text) == 0:
                                        continue
                            else:
                                # filter
                                continue

                            ann_box.append([xc, yc, shape.width, shape.height])
                            ann_cat.append(self.label2index[cur_class])
                        except:
                            pass

                    if len(ann_cat) > max_element_num or len(ann_cat) <= 3:
                        continue

                    ann_box = np.array(ann_box)
                    ind = np.lexsort((ann_box[:, 0], ann_box[:, 1]))
                    # Sort by ann_box[:, 1], then by ann_box[:, 0]
                    ann_box = ann_box[ind]
                    
                    ann_cat = np.array(ann_cat)
                    
                    ann_cat = ann_cat[ind]

                    ann_box[:, [0, 2]] = ann_box[:, [0, 2]]/(W -1)
                    ann_box[:, [1, 3]] = ann_box[:, [1, 3]]/(H -1)

                    boxes = torch.tensor(ann_box, dtype=torch.float)
                    labels = torch.tensor(ann_cat, dtype=torch.long)

                    ppt_name = file.split(".")[0]
                    data = Data(x=boxes, y=labels)
                    data.attr = {
                        'name':f"{ppt_name}_{index}",
                        'width': W,
                        'height': H,
                        'has_canvas_element': False,
                    }
                    data_list.append(data)
                        
        # train 85% / val 5% / test 10%
        N = len(data_list)
        s = [int(N * .85), int(N * .90)]
        torch.save(self.collate(data_list[:s[0]]), self.processed_paths[0])
        torch.save(self.collate(data_list[s[0]:s[1]]), self.processed_paths[1])
        torch.save(self.collate(data_list[s[1]:]), self.processed_paths[2])